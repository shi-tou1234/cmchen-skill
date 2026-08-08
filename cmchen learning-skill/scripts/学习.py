#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
学习 Skill — 内容提取工具
把 PDF / PPTX / PPT 里的文字提取成 JSON，供交互式学习使用。
用法：
    python 学习.py extract <目录> -o <输出.json>
"""
import os, sys, json, re
sys.stdout.reconfigure(encoding='utf-8')


def extract_content(input_dir):
    files = []
    for root, dirs, fnames in os.walk(input_dir):
        for f in fnames:
            if f.endswith((".pdf", ".pptx", ".ppt")):
                files.append(os.path.join(root, f))
    if not files:
        print("  [!] 目录里没找到 .pdf/.pptx/.ppt 文件")
        return []
    results = []
    for fp in sorted(files):
        rel = os.path.relpath(fp, input_dir)
        ext = os.path.splitext(fp)[1].lower()
        text = None
        try:
            if ext == ".pdf":
                text = _extract_pdf(fp)
            elif ext == ".pptx":
                text = _extract_pptx(fp)
            elif ext == ".ppt":
                text = _extract_ppt_clean(fp)
        except Exception as e:
            print(f"  [!] Error: {rel} -> {e}")
            continue
        if text and len(text) > 50 and not _is_garbage(text):
            results.append({"file": rel, "text": text, "type": ext})
            print(f"  [OK] {rel} ({len(text)} chars)")
        else:
            print(f"  [--] {rel} (太短或乱码，跳过)")
    return results


def _is_garbage(text):
    valid = sum(1 for c in text if "一" <= c <= "鿿" or " " <= c <= "~")
    return valid / max(1, len(text)) < 0.3


def _extract_pdf(fp):
    import PyPDF2
    return "\n\n".join(
        p.extract_text().strip()
        for p in PyPDF2.PdfReader(fp).pages
        if p.extract_text().strip()
    )


def _extract_pptx(fp):
    from pptx import Presentation
    pages = []
    for slide in Presentation(fp).slides:
        texts = [
            p.text.strip()
            for s in slide.shapes
            if s.has_text_frame
            for p in s.text_frame.paragraphs
            if p.text.strip()
        ]
        if texts:
            pages.append("\n".join(texts))
    return "\n\n---\n\n".join(pages)


def _extract_ppt_clean(fp):
    import olefile
    ole = olefile.OleFileIO(fp)
    data = ole.openstream("PowerPoint Document").read()
    ole.close()
    chars = []
    i = 0
    while i < len(data) - 1 and len(chars) < 50000:
        low, high = data[i], data[i + 1]
        if high == 0 and 0x20 <= low < 0x7f:
            chars.append("\n" if low == 0x0d else chr(low))
        elif 0x4E <= high <= 0x9F:
            try:
                chars.append(data[i : i + 2].decode("utf-16-le"))
            except Exception:
                pass
        i += 2
    return _clean_ppt_text("".join(chars))


def _clean_ppt_text(text, max_len=50000):
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", text)
    text = re.sub(r"(.)\\1{10,}", "", text)
    cleaned = []
    for c in text:
        if "一" <= c <= "鿿" or " " <= c <= "~" or c in "\n\r":
            cleaned.append(c)
        elif cleaned and cleaned[-1] not in " \n":
            cleaned.append(" ")
    result = "".join(cleaned)
    lines = []
    for line in result.split("\n"):
        line = line.strip()
        if len(line) < 3:
            continue
        cn = sum(1 for c in line if "一" <= c <= "鿿")
        if cn == 0 and len(line) < 10:
            continue
        lines.append(line)
    return "\n".join(lines)


def cmd_extract(args):
    data = extract_content(args.input_dir)
    with open(args.output, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    print(f"[OK] {len(data)} files -> {args.output}")


if __name__ == "__main__":
    import argparse
    p = argparse.ArgumentParser(description="学习 Skill — 内容提取")
    sub = p.add_subparsers(dest="cmd")
    sp = sub.add_parser("extract", help="从目录提取文本（pdf/pptx/ppt）")
    sp.add_argument("input_dir", type=str)
    sp.add_argument("--output", "-o", type=str, default="content.json")
    a = p.parse_args()
    if a.cmd == "extract":
        cmd_extract(a)
    else:
        p.print_help()
